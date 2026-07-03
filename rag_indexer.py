"""
rag_indexer.py — クライアント資料のRAG索引を作成する

使い方:
  py rag_indexer.py --client 田所商店

索引ファイル: ../rag_index_{client}.json
"""

import argparse
import json
import logging
import os
import sys
import time
import difflib
from pathlib import Path

os.environ.setdefault("HF_HUB_DISABLE_PROGRESS_BARS", "1")
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
logging.getLogger("transformers").setLevel(logging.ERROR)
logging.getLogger("sentence_transformers").setLevel(logging.ERROR)
logging.getLogger("huggingface_hub").setLevel(logging.ERROR)

BASE_DIR   = Path(__file__).parent.parent  # Documents/gijiroku/
OUTPUT_DIR = BASE_DIR / "output"
ALIASES_FILE = BASE_DIR / "client_aliases.json"
FOLDERS_FILE = BASE_DIR / "client-folders.json"

CHUNK_SIZE    = 600   # 目標チャンクサイズ（文字数）
CHUNK_OVERLAP = 100   # 前後のオーバーラップ（文字数）
ALIAS_THRESHOLD = 0.75  # エイリアス一致の類似度閾値


def load_aliases():
    with open(ALIASES_FILE, encoding="utf-8") as f:
        return json.load(f)


def normalize_client_name(raw_name: str, aliases: dict) -> str | None:
    """
    rawをエイリアス対応表と照合し正式名を返す。
    一致なしはNone。
    """
    for official, alias_list in aliases.items():
        for alias in alias_list:
            ratio = difflib.SequenceMatcher(None, raw_name, alias).ratio()
            if ratio >= ALIAS_THRESHOLD:
                return official
    return None


def scan_output_dir(client_name: str, aliases: dict):
    """
    output/ を走査しクライアント判定。
    ファイル名規則: 議事録_{クライアント名}_{会議名}_{日付}.docx
                   または 議事録_{クライアント名}_{日付}.docx
    """
    matched = []
    unclassified = []

    for f in OUTPUT_DIR.iterdir():
        if not f.is_file():
            continue
        parts = f.stem.split("_")
        if len(parts) < 2 or parts[0] != "議事録":
            unclassified.append(str(f))
            continue

        raw = parts[1]
        normalized = normalize_client_name(raw, aliases)
        if normalized == client_name:
            matched.append(f)
        elif normalized is None:
            # どのクライアントにも一致しない
            if not any(normalize_client_name(raw, aliases) == n for n in aliases):
                unclassified.append(str(f))

    return matched, unclassified


# ─── テキスト抽出 ────────────────────────────────────────────

def extract_text_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text_pdf(path: Path) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n".join(texts)


def extract_text_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            line = "\t".join(str(v) for v in row if v is not None)
            if line.strip():
                texts.append(line)
    return "\n".join(texts)


def extract_text(path: Path) -> str | None:
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            return extract_text_docx(path)
        elif ext == ".pptx":
            return extract_text_pptx(path)
        elif ext == ".pdf":
            return extract_text_pdf(path)
        elif ext == ".xlsx":
            return extract_text_xlsx(path)
        elif ext in (".txt", ".md"):
            return path.read_text(encoding="utf-8", errors="ignore")
        else:
            return None  # 非対応形式
    except Exception as e:
        print(f"  [WARN] テキスト抽出失敗 {path.name}: {e}", file=sys.stderr)
        return None


# ─── チャンク分割 ─────────────────────────────────────────────

def chunk_text(text: str, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - overlap
    return [c for c in chunks if c.strip()]


# ─── 埋め込み ─────────────────────────────────────────────────

_model = None

def get_model():
    global _model
    if _model is None:
        from sentence_transformers import SentenceTransformer
        print("埋め込みモデルをロード中（初回はダウンロードあり）...", file=sys.stderr)
        _model = SentenceTransformer("intfloat/multilingual-e5-small")
    return _model


def embed_passages(texts: list[str]) -> list[list[float]]:
    model = get_model()
    # multilingual-e5 は passage: プレフィックスが必要
    prefixed = ["passage: " + t for t in texts]
    vecs = model.encode(prefixed, normalize_embeddings=True, batch_size=32, show_progress_bar=False)
    return [v.tolist() for v in vecs]


# ─── メイン ───────────────────────────────────────────────────

def collect_files(client_name: str):
    """顧客フォルダ + 議事録フォルダから対象ファイルを収集"""
    aliases = load_aliases()

    # 顧客フォルダ
    with open(FOLDERS_FILE, encoding="utf-8") as f:
        folders = json.load(f)
    client_folder = next((e["path"] for e in folders if e["client"] == client_name), None)

    client_files = []  # (path, source_type)
    if client_folder:
        p = Path(client_folder)
        if p.exists():
            for f in p.rglob("*"):
                if f.is_file():
                    client_files.append((f, "顧客フォルダ"))
        else:
            print(f"[WARN] 顧客フォルダが存在しません: {client_folder}", file=sys.stderr)
    else:
        print(f"[WARN] client-folders.json に {client_name} の登録なし", file=sys.stderr)

    # 議事録フォルダ
    matched, unclassified = scan_output_dir(client_name, aliases)
    for f in matched:
        client_files.append((f, "議事録フォルダ"))

    if unclassified:
        print(f"[未分類ファイル] {len(unclassified)} 件:", file=sys.stderr)
        for u in unclassified:
            print(f"  {u}", file=sys.stderr)

    return client_files, len(matched), len(unclassified)


def build_index(client_name: str):
    t0 = time.time()
    print(f"=== RAG索引作成開始: {client_name} ===")

    files, gijiroku_count, unclassified_count = collect_files(client_name)
    print(f"対象ファイル: {len(files)} 件（議事録 {gijiroku_count} 件、未分類 {unclassified_count} 件）")

    all_chunks = []
    processed = 0
    skipped = 0

    for path, source_type in files:
        text = extract_text(path)
        if text is None:
            print(f"  [SKIP] {path.name} （非対応形式）")
            skipped += 1
            continue
        if not text.strip():
            print(f"  [SKIP] {path.name} （テキストなし）")
            skipped += 1
            continue

        chunks = chunk_text(text)
        for i, chunk in enumerate(chunks):
            all_chunks.append({
                "text":        chunk,
                "file_name":   path.name,
                "file_path":   str(path),
                "source_type": source_type,
                "chunk_index": i,
            })
        print(f"  [OK]   {path.name} — {len(chunks)} チャンク ({source_type})")
        processed += 1

    if not all_chunks:
        print("有効なチャンクがありません。索引を作成できません。")
        return

    # 埋め込み
    print(f"\n{len(all_chunks)} チャンクを埋め込み中...")
    vectors = embed_passages([c["text"] for c in all_chunks])
    for chunk, vec in zip(all_chunks, vectors):
        chunk["embedding"] = vec

    # 保存
    index_path = BASE_DIR / f"rag_index_{client_name}.json"
    index_data = {
        "client":    client_name,
        "built_at":  time.strftime("%Y-%m-%dT%H:%M:%S"),
        "chunks":    all_chunks,
    }
    with open(index_path, "w", encoding="utf-8") as f:
        json.dump(index_data, f, ensure_ascii=False)

    elapsed = time.time() - t0
    print(f"\n=== 完了 ({elapsed:.1f}秒) ===")
    print(json.dumps({
        "status":          "ok",
        "client":          client_name,
        "processed_files": processed,
        "skipped_files":   skipped,
        "unclassified":    unclassified_count,
        "total_chunks":    len(all_chunks),
        "index_path":      str(index_path),
        "elapsed_sec":     round(elapsed, 1),
    }, ensure_ascii=False))


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--client", required=True, help="クライアント名")
    args = parser.parse_args()
    build_index(args.client)
